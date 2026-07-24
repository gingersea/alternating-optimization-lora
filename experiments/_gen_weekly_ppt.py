"""Generate weekly group meeting PPT: A-SYNC protocol evolution + FLOPs experiment.

14 slides covering this week's 5 commits (Jul 23-24, 2026).
Design: 16:9 white bg, dark-blue title bar, figure-heavy, CN text.
"""

import os, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W, H = Inches(13.333), Inches(7.5)
DARK = RGBColor(0x1F, 0x49, 0x7D)
RED = RGBColor(0xC0, 0x50, 0x4D)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x9B, 0xBB, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)
GRAY = RGBColor(0x60, 0x60, 0x60)
BLACK = RGBColor(0x20, 0x20, 0x20)

FIG_DIR = os.path.join(os.path.dirname(__file__), "..", "docs", "figures")
FINAL = os.path.join(FIG_DIR, "final_report")
BILINGUAL = os.path.join(FIG_DIR, "bilingual_report")

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]  # blank


def title_bar(slide, text, subtitle=None):
    """Add dark-blue title bar at top."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p2.font.italic = True


def add_text(slide, left, top, width, height, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=14, color=BLACK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
        p.level = 0
    return tf


def add_img(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        add_text(slide, left, top, 3, 0.5, f"[Missing: {os.path.basename(path)}]", size=10, color=RED)
        return
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    slide.shapes.add_picture(path, Inches(left), Inches(top), **kw)


def section_slide(text, subtitle=""):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()
    add_text(slide, 1, 3.5, 11, 1.5, text, size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, 1, 4.5, 11, 1, subtitle, size=16, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
    return slide


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
add_text(s, 1, 1.5, 11, 1.5, "A-SYNC 协议族演化\n与 FLOPs 归一化对比实验", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1, 4, 11, 1, "Alternating Optimization for LLM Post-Training — 周进展汇报", size=18, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
add_text(s, 1, 5.2, 11, 0.6, "2026年7月23-24日  |  5 commits  |  2 GPU实验  |  3份文档", size=14, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: Overview — 本周做了什么
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "本周工作概览", "5 commits, 3 major deliverables")
add_bullets(s, 0.6, 1.4, 12, 5.5, [
    "📊  双语文档交付：EN + ZH A-SYNC 变体报告（456KB + 529KB PDF）",
    "📐  基础理论文档：残差放大（ρ≈1.08）因果推导 + A-SYNC 12变体分类",
    "🧪  FLOPs归一化实验：A-SYNC CONSTANT vs AdamW vs LoRA on OPT-125m",
    "",
    "关键数字：",
    "  • 5 commits 推送至 github.com/hjiang555-a11y + gingersea",
    "  • 7 张新图表 + 1 张 FLOPs 对比图",
    "  • ~4000 字技术文档（中英双语）",
    "  • A-SYNC CONSTANT 48c: OPT-125m PPL 2246→60.7, Qwen7B PPL 58.8→7.6",
], size=15)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: Background — 问题根源
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "问题背景：残差连接 = 扰动放大器", "为什么 Protocol A 在深层模型上发散？")
add_bullets(s, 0.6, 1.3, 5.8, 5.5, [
    "每层 Transformer: h_{l+1} = h_l + f_l(h_l)",
    "恒等路径 h_l 原样保留 ALS 扰动 δ",
    "每层额外叠加 f_l 对扰动输入的响应",
    "",
    "核心公式：δ_{l+1} = (I + J_l) · δ_l",
    "  J_l = 层 l 的雅可比矩阵",
    "  ‖I + J_l‖ ≈ 1.08 (跨层几何均值)",
    "",
    "级联 27 层: 1.08²⁷ ≈ 8.0× 放大",
    "SGD 每周期恢复 ≈ 0.005 PPL",
    "→ 1600:1 不对称 → 灾难性发散",
], size=13)
add_img(s, os.path.join(FINAL, "fig1_residual.png"), 7, 1.2, width=5.8, height=3.2)
add_img(s, os.path.join(BILINGUAL, "fig6_depth_comparison.png"), 7, 4.5, width=5.8, height=2.8)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: A-SYNC mechanism
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "核心创新：A-SYNC 梯度注入机制", "从权重修改 → 梯度塑形")
add_bullets(s, 0.6, 1.3, 5.5, 5.5, [
    "Protocol A (旧): ALS 直接修改 lm_head",
    "  → 权重变化进入前传 → 残差链放大",
    "  → 28层模型 11/11 次发散",
    "",
    "A-SYNC (新): 3步流程",
    "  ① ALS 计算 δ = W_als − W_before",
    "  ② 权重立即恢复！lm_head 不变",
    "  ③ SGD 时注入: grad += sync × δ",
    "",
    "关键：模型前传从不看到 ALS 权重",
    "  → 扰动绕过残差放大路径",
    "  → 7B (28L) 首次收敛！",
], size=13)
add_img(s, os.path.join(FINAL, "fig2_arch.png"), 6.5, 1.2, width=6.3, height=2.8)
add_img(s, os.path.join(FINAL, "fig4_convergence.png"), 6.5, 4.2, width=6.3, height=3.0)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: Variant table
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "A-SYNC 变体演化全景", "12 variants → CONSTANT 48c 最优 (7B PPL=7.6)")
rows = [
    ("变体", "机制要点", "7B PPL", "vs CONSTANT"),
    ("A-SYNC Vanilla (8c)", "delta注入+扰动+指数衰减", "25.8", "Δ+18.2"),
    ("A-SYNC No-Perturb (8c)", "去扰动阶段", "16.6", "Δ+9.0"),
    ("A-SYNC Cosine (32c)", "余弦衰减 sync+lr", "13.2", "Δ+5.6"),
    ("A-SYNC CONSTANT (24c)", "恒sync=0.05, lr=2e-4", "9.0", "Δ+1.4"),
    ("A-SYNC CONSTANT (48c)", "延至48周期, C44收敛", "7.6 ★", "BASELINE"),
    ("A-CYCLE (3×8)", "Cosine重启×3块", "16.5", "Δ+8.9"),
    ("A-SYNC+EMA", "δ的指数滑动平均", "0.5B: 5.5", "未测7B"),
    ("A-SYNC+Aligned", "只注入梯度同向分量", "0.5B: 5.5", "未测7B"),
    ("A-SYNC+SWA", "权重平均 (C10起)", "10.5", "Δ+2.9"),
    ("A-PROBE (r=64)", "低秩探针绕过lm_head", "22.8", "Δ+15.2"),
    ("LARS optimizer", "层自适应学习率", "未测7B", "—"),
]
y0 = 1.3
col_x = [0.3, 3.3, 6.8, 10.0, 11.5]
col_w = [2.9, 3.4, 3.0, 2.0, 1.5]
for r, row in enumerate(rows):
    for c, (text, x, w) in enumerate(zip(row, col_x, col_w)):
        clr = WHITE if r == 0 else (LIGHT_BG if r % 2 == 0 else WHITE)
        fc = WHITE if r == 0 else BLACK
        sz = 11 if r == 0 else 10
        bd = True if r == 0 else False
        is_best = "★" in text or (r == 5 and c == 1)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y0 + r * 0.42),
                                 Inches(w), Inches(0.4))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK if r == 0 else (RGBColor(0xE8, 0xF0, 0xF8) if is_best else clr)
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_top = Inches(0.02)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = fc
        p.font.bold = bd or is_best

add_text(s, 0.6, y0 + len(rows) * 0.42 + 0.15, 12, 0.4,
         "结论: 衰减策略(cosine/exponential)全部有害 → CONSTANT 无衰减是最优调度", size=11, color=RED, bold=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: 7B full convergence table
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "7B 收敛排名：全部 A-SYNC 变体", "Qwen2.5-7B (28L) — CONSTANT 48c 夺冠")
add_img(s, os.path.join(FINAL, "fig7_scoreboard.png"), 0.3, 1.2, width=12.5, height=5.8)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: Fix taxonomy
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "修复策略分类：5大类别", "哪个策略真正有效？")
add_bullets(s, 0.6, 1.3, 12, 5.5, [
    "A. 降低扰动幅度 — EMA平滑、Aligned注入、step_size减小 → 效果有限",
    "B. 增强恢复能力 — 更多SGD步、更高lr、LARS → 边际改善",
    "C. 改变干预机制 ★ — A-SYNC梯度注入(关键突破)、A-PROBE低秩探针 → 一阶改善",
    "D. 改变衰减调度 ★★ — CONSTANT无衰减(最优)、Cosine(最差)、Warm Restart(中等)",
    "E. 事后平滑 — SWA → 在A-SYNC非稳态轨迹上反而变差(+3.3 PPL)",
    "",
    "核心洞察:",
    "  • Category C (梯度注入) 是让7B收敛的唯一原因",
    "  • Category D (CONSTANT) 是7B PPL从25.8→7.6的关键",
    "  • ALS δ幅值随body自适应自动衰减 → 外部衰减完全多余",
], size=13)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: Section divider — FLOPs experiment
# ═══════════════════════════════════════════════════════════════════
section_slide("FLOPs 归一化对比实验", "A-SYNC CONSTANT vs AdamW Full-Rank vs LoRA AdamW  on  OPT-125m")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: FLOPs experiment setup
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "FLOPs 归一化实验设计", "OPT-125m (12L, 125M params), WikiText-2, float32")
add_bullets(s, 0.6, 1.3, 5.5, 5.5, [
    "3协议对比:",
    "  ① A-SYNC CONSTANT 48c",
    "     ALS δ → grad注入, sync=0.05恒",
    "  ② AdamW Full-Rank (Protocol B)",
    "     β=(0.9,0.999), lr=1e-4, wd=0.01",
    "  ③ LoRA AdamW r=8, α=16",
    "     target: q,v,k,out_proj",
    "",
    "FLOPs 会计 (每步):",
    "  ALS = 4 × params = 501 MFLOPs",
    "  SGD = 6 × params = 751 MFLOPs",
    "  AdamW = 10 × params = 1252 MFLOPs",
    "  LoRA AdamW = 10 × lora = 5.9 MFLOPs",
    "",
    "FLOPs 预算匹配: AdamW步数 = A-SYNC总FLOPs/10",
    "GPU: 2×RTX 5090, ~5min/协议",
], size=12)

add_bullets(s, 6.5, 1.3, 6.3, 5.5, [
    "关键参数:",
    "  • 125,239,296 全秩参数",
    "  • 589,824 LoRA参数 (213× fewer)",
    "  • 24 + 48 周期 A-SYNC",
    "  • 720 步 AdamW/LoRA",
    "  • batch_size=4 train / 8 eval",
    "  • seq_len=128",
    "",
    "评估: 每个 A-SYNC cycle / 每30 AdamW步",
    "  → 25 个评估点/协议",
    "",
    "PPL计算: WikiText-2 test set",
    "  PPL = exp(avg cross-entropy)",
], size=12)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: FLOPs main plot
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "FLOPs vs PPL 对比图", "OPT-125m: AdamW 最优 (PPL 23.2), LoRA 最省 (0.013T)")
add_img(s, os.path.join(FIG_DIR, "flops_sweep_opt125m.png"), 0.3, 1.15, width=12.8, height=6.2)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: FLOPs results table
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "FLOPs 归一化对比数据", "数值结果")
data_rows = [
    ("协议", "最终 PPL", "FLOPs (T)", "墙钟时间", "PPL/TFLOP"),
    ("AdamW Full-Rank",   "23.2", "0.911",  "136s", "25.5"),
    ("LoRA AdamW (r=8)",  "37.3", "0.013",  "142s", "2812.4"),
    ("A-SYNC CONSTANT 48c","60.7","1.846",  "318s", "32.9"),
    ("A-SYNC CONSTANT 24c","74.1","0.923",  "163s", "80.3"),
]
y0 = 1.3
col_x = [0.5, 3.3, 6.5, 9.0, 11.0]
col_w = [2.7, 3.0, 2.3, 2.5, 2.0]
for r, row in enumerate(data_rows):
    for c, (text, x, w) in enumerate(zip(row, col_x, col_w)):
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y0 + r * 0.5),
                                 Inches(w), Inches(0.48))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK if r == 0 else (RGBColor(0xE8, 0xF0, 0xF8) if r == 1 else (LIGHT_BG if r % 2 == 0 else WHITE))
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12 if r == 0 else 11)
        p.font.color.rgb = WHITE if r == 0 else BLACK
        p.font.bold = (r == 0 or r == 1)

add_bullets(s, 0.6, y0 + len(data_rows) * 0.5 + 0.2, 12, 2.5, [
    "• AdamW Full-Rank 绝对最优: PPL 23.2 at 0.91T FLOPs",
    "• LoRA AdamW 效率最优: PPL 37.3 at 0.013T FLOPs — 70× 更省算力",
    "• A-SYNC 48c 仍在收敛: 每cycle ΔPPL ≈ −0.3~−0.7, 未见平台 (预测80+cycle可达PPL<40)",
    "• A-SYNC 优势在深层模型(28L+)而非浅层(12L): ρ^11=2.3× ← 梯度注入无收益",
], size=12)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: Trajectory detail
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "A-SYNC 48c 收敛轨迹 (OPT-125m)", "单调收敛，未见平台")
add_bullets(s, 0.6, 1.3, 5.5, 5.5, [
    "起始: PPL = 2246.1 (baseline)",
    "Cycle  1: 1458.4   Cycle 13:  97.7",
    "Cycle  5:  360.2   Cycle 17:  84.6",
    "Cycle  9:  140.2   Cycle 21:  78.1",
    "",
    "Cycle 25:  73.3   Cycle 37:  65.0",
    "Cycle 29:  69.9   Cycle 41:  63.2",
    "Cycle 33:  67.2   Cycle 45:  61.8",
    "                    Cycle 48:  60.7",
    "",
    "晚期每cycle改善: ΔPPL ≈ −0.3~−0.7",
    "预测: 100 cycle → PPL ≈ 40",
    "AdamW 在相同 FLOPs: PPL = 23.2 (已趋平)",
    "",
    "结论: A-SYNC 持续改善但速度慢",
    "  → 需 3-4× AdamW算力才能追平",
], size=12)

add_img(s, os.path.join(FIG_DIR, "flops_sweep_opt125m.png"), 6.5, 1.3, width=6.3, height=5.5)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13: Cross-model perspective
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "跨模型结果汇总", "不同深度的不同表现")
add_bullets(s, 0.6, 1.3, 5.5, 5, [
    "12层 (OPT-125m):",
    "  • ρ^11 = 2.3× → A-SYNC收敛但慢",
    "  • AdamW 大胜: PPL 23.2 vs 60.7",
    "",
    "24层 (Qwen0.5B):",
    "  • ρ^23 = 5.9× → 边际不稳定",
    "  • A-SYNC+ 全部收敛 PPL 5.5 (容量底)",
    "  • LoRA 同达 PPL 5.5 (更省算力)",
    "",
    "28层 (Qwen7B):",
    "  • ρ^27 = 8.0× → Protocol A 发散",
    "  • A-SYNC CONSTANT 首次收敛: PPL 7.6",
    "  • AdamW 仍远优: PPL 1.25 (800 steps)",
    "  • 差距 6.1× — 但 7.7× 低于 baseline",
], size=12)

add_bullets(s, 6.5, 1.3, 6.3, 5, [
    "关键洞察:",
    "",
    "1. A-SYNC 价值在深层模型(≥28L)",
    "   浅层模型梯度注入开销无收益",
    "",
    "2. CONSTANT 调度是核心发现",
    "   衰减策略全部反作用: Cosine 最差,",
    "   Exponential 次之, 无衰减最优",
    "",
    "3. 自然收敛是确证的",
    "   7B C44 收敛, OPT-125m 48c 未见平台",
    "   → ALS δ 幅值随 body 自适应自动衰减",
    "   → 外部衰减完全多余",
    "",
    "4. 仍有 6.1× 差距 vs AdamW on 7B",
    "   可能原因: ALS 目标不对齐 + 仅lm_head",
    "   未覆盖 + 梯度裁剪限幅",
], size=12)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14: Next steps + conclusion
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "下一步计划 & 总结", "What's next?")

add_bullets(s, 0.6, 1.3, 5.8, 5.5, [
    "短期 (1-2周):",
    "  1. A-SYNC+EMA/Aligned 7B 验证",
    "     (脚本已有, 仅缺GPU运行)",
    "  2. A-PROBE 更大rank (256/512/1024)",
    "  3. 多head ALS: 2-3个attention层",
    "",
    "中期 (2-4周):",
    "  4. A-SYNC CONSTANT 96+ cycle 7B",
    "     (测试是否能追平AdamW)",
    "  5. A-SYNC+LARS 7B (层自适应学习率)",
    "",
    "论文推进:",
    "  6. 独立复核者复算主表",
    "  7. 选择投稿期刊 (TMLR / arXiv+workshop)",
], size=13)

add_bullets(s, 6.8, 1.3, 6, 5.5, [
    "本周核心结论:",
    "",
    "1. 残差放大 (ρ≈1.08) 是Protocol A",
    "   发散的根本原因 (因果推导+实验验证)",
    "",
    "2. A-SYNC 梯度注入绕过残差放大",
    "   → 7B首次收敛 (从发散 → PPL 7.6)",
    "",
    "3. CONSTANT 调度是最优的:",
    "   衰减全部反作用, 自然收敛无需外部衰减",
    "",
    "4. FLOPs归一化: AdamW > A-SYNC",
    "   但 A-SYNC 仍在改善 (未见平台)",
    "   LoRA 效率领先 70×",
    "",
    "5. A-SYNC 定位: 深层模型专用",
    "   12L 无优势, 28L 是关键突破",
], size=13, color=DARK)

# thank-you bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), W, Inches(0.5))
bar.fill.solid(); bar.fill.fore_color.rgb = DARK; bar.line.fill.background()
tf = bar.text_frame
p = tf.paragraphs[0]
p.text = "Thanks!  ·  github.com/hjiang555-a11y/alternating-optimization-lora"
p.font.size = Pt(11)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ── Save ────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "..", "docs", "weekly_group_meeting_20260724.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
